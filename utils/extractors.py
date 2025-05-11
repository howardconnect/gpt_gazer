import os
import fitz  # PyMuPDF
from utils.ocr import ocr_image_from_pdf
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
from email import policy
from email.parser import BytesParser
from logger import logger  # ✅ added centralized logging

def extract_text_from_file(filepath):
    ext = os.path.splitext(filepath)[1].lower()

    try:
        if ext in [".txt", ".md", ".log"]:
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()

        elif ext == ".html":
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                soup = BeautifulSoup(f, "html.parser")
                return soup.get_text(separator="\n")

        elif ext == ".eml":
            with open(filepath, "rb") as f:
                msg = BytesParser(policy=policy.default).parse(f)
                body = msg.get_body(preferencelist=('plain'))
                return "\n".join([
                    f"From: {msg['From']}",
                    f"To: {msg['To']}",
                    f"Subject: {msg['Subject']}",
                    f"Date: {msg['Date']}",
                    f"\n{body.get_content() if body else '(No plain text body)'}"
                ])

        elif ext == ".pdf":
            doc = fitz.open(filepath)
            text = "\n".join(page.get_text() for page in doc)
            if text.strip():
                return text
            logger.warning(f"📄 No text found in PDF — falling back to OCR: {filepath}")
            return ocr_image_from_pdf(filepath)

        elif ext == ".docx":
            doc = Document(filepath)
            return "\n".join(p.text for p in doc.paragraphs)

        elif ext == ".pptx":
            prs = Presentation(filepath)
            return "\n".join(
                shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
            )

        elif ext in [".json", ".xml", ".rtf"]:
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()

        else:
            logger.warning(f"⚠️ Unsupported file type for text extraction: {ext}")
            return ""

    except Exception as e:
        logger.exception(f"❌ Extraction failed for {filepath}: {e}")
        return ""
